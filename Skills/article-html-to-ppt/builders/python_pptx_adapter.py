from __future__ import annotations

import importlib.metadata
from pathlib import Path
from typing import Any

from .base import BuildInspection, BuildPlan, BuildResult, BuilderCapability, SupportLevel


class PythonPptxAdapter:
    name = "python_pptx"

    def probe(self) -> BuilderCapability:
        version = _package_version()
        available = version is not None
        return BuilderCapability(
            name=self.name,
            available=available,
            version=version,
            command="python" if available else None,
            features={
                "native_text": True,
                "native_table": True,
                "native_chart": True,
                "native_connector": True,
                "svg_embed": False,
                "speaker_notes": False,
                "theme": "partial",
                "slide_master": "partial",
                "readback": True,
            },
            warnings=[] if available else ["python-pptx package is not importable."],
        )

    def supports(self, component_type: str, delivery_route: str) -> SupportLevel:
        if delivery_route in {"native_ppt", "native_table", "native_chart", "native_diagram"}:
            return "full"
        if delivery_route in {"generated_image", "background_image", "hybrid_overlay"}:
            return "visual_only"
        return "unsupported"

    def plan(self, ppt_ir: dict[str, Any], style_contract: dict[str, Any], delivery_plan: dict[str, Any]) -> BuildPlan:
        object_plans = {
            (obj.get("slide_id"), obj.get("object_id")): obj
            for slide in delivery_plan.get("slides", []) or []
            for obj in slide.get("objects", []) or []
            if isinstance(obj, dict)
        }
        slides = []
        for slide in ppt_ir.get("slides", []) or []:
            planned_objects = []
            for obj in slide.get("objects", []) or []:
                planned = dict(obj)
                planned["delivery_plan"] = object_plans.get((slide.get("id"), obj.get("id")), {})
                planned_objects.append(planned)
            slide_plan = dict(slide)
            slide_plan["objects"] = planned_objects
            slides.append(slide_plan)
        warnings = _unsupported_style_warnings(style_contract)
        return BuildPlan(builder=self.name, slides=slides, style_contract=dict(style_contract), warnings=warnings)

    def build(self, build_plan: BuildPlan, output_dir: Path) -> BuildResult:
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
            from pptx.util import Inches, Pt
        except Exception as exc:
            return BuildResult(builder=self.name, status="failed", errors=[{"code": "BUILDER_RUNTIME_IMPORT_FAILED", "message": str(exc)}])

        output_dir.mkdir(parents=True, exist_ok=True)
        deck_path = output_dir / "deck.pptx"
        style = _resolve_style(build_plan.style_contract)
        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        blank_layout = presentation.slide_layouts[6]

        object_results: list[dict[str, Any]] = []
        fallbacks: list[dict[str, Any]] = []
        warnings: list[str] = list(build_plan.warnings)

        for slide_plan in build_plan.slides:
            slide = presentation.slides.add_slide(blank_layout)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor.from_string(style["background"])
            _add_textbox(
                slide,
                slide_plan.get("title") or "Untitled",
                x=style["margin_left"],
                y=style["margin_top"],
                w=12.2,
                h=0.45,
                font_size=style["title_size"],
                bold=True,
                color=style["primary"],
                font_name=style["font_name"],
            )
            body_parts = [part for part in [slide_plan.get("judgment"), slide_plan.get("message")] if isinstance(part, str) and part.strip()]
            body = "\n".join(body_parts)
            if body:
                _add_textbox(slide, body, x=style["margin_left"], y=1.05, w=11.7, h=0.5, font_size=style["body_size"], color=style["text_secondary"], font_name=style["font_name"])
            objects = slide_plan.get("objects", []) or []
            if not objects:
                continue
            card_width = max(2.6, min(3.6, 10.8 / max(len(objects), 1)))
            for index, obj in enumerate(objects):
                route = ((obj.get("delivery_plan") or {}).get("selected_route") or obj.get("delivery_preferences", {}).get("preferred_route") or "native_ppt")
                actual_route = route
                component_type = obj.get("component_type") or obj.get("type")
                is_chart = route == "native_chart" and obj.get("type") == "chart"
                is_process = route == "native_diagram" and component_type == "process"
                if route in {"hybrid_overlay", "svg_component", "generated_image", "background_image"}:
                    actual_route = "native_ppt"
                    fallbacks.append(
                        {
                            "slide_id": slide_plan.get("id"),
                            "object_id": obj.get("id"),
                            "component_type": component_type,
                            "planned_route": route,
                            "actual_route": actual_route,
                            "reason_codes": ["PYTHON_PPTX_MINIMAL_NATIVE_FALLBACK"],
                            "editable_core_preserved": obj.get("editability") != "native_required",
                        }
                    )
                x = style["margin_left"] + index * (card_width + style["card_gap"])
                y = 1.85
                if route == "native_table" or obj.get("type") == "table":
                    _add_table(slide, obj, x=x, y=y, w=min(10.8, card_width * 2), h=2.3, style=style)
                elif is_chart:
                    _add_chart(slide, obj, x=x, y=y, w=card_width, h=3.4, style=style)
                elif is_process:
                    _add_process(slide, obj, x=x, y=y, w=card_width, h=1.35, style=style)
                else:
                    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_width), Inches(1.35))
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor.from_string(style["surface"])
                    shape.line.color.rgb = RGBColor.from_string(style["border"])
                    text_frame = shape.text_frame
                    text_frame.clear()
                    paragraph = text_frame.paragraphs[0]
                    run = paragraph.add_run()
                    run.text = _object_text(obj)
                    run.font.size = Pt(style["body_size"])
                    run.font.name = style["font_name"]
                    run.font.color.rgb = RGBColor.from_string(style["text_primary"])
                object_results.append(
                    {
                        "slide_id": slide_plan.get("id"),
                        "object_id": obj.get("id"),
                        "component_type": obj.get("component_type") or obj.get("type"),
                        "planned_route": route,
                        "actual_route": actual_route,
                        "status": "created",
                    }
                )

        presentation.save(deck_path)
        return BuildResult(builder=self.name, status="created", pptx=str(deck_path), object_results=object_results, fallbacks=fallbacks, warnings=warnings)

    def inspect_output(self, pptx_path: Path) -> BuildInspection:
        try:
            from ppt_qa.package_inspector import inspect_package
        except Exception as exc:
            return BuildInspection(builder=self.name, status="failed", errors=[{"code": "BUILDER_INSPECT_IMPORT_FAILED", "message": str(exc)}])
        inspection = inspect_package(pptx_path)
        return BuildInspection(
            builder=self.name,
            status="passed" if inspection.status == "passed" else "failed",
            issues=[issue.__dict__ for issue in inspection.issues],
        )


def _package_version() -> str | None:
    try:
        return importlib.metadata.version("python-pptx")
    except importlib.metadata.PackageNotFoundError:
        return None


def _add_textbox(slide: Any, text: str, *, x: float, y: float, w: float, h: float, font_size: float, color: str, bold: bool = False, font_name: str | None = None) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = RGBColor.from_string(color)


def _object_text(obj: dict[str, Any]) -> str:
    content = obj.get("content")
    if isinstance(content, str) and content.strip():
        return content.strip()
    if isinstance(content, dict):
        for key in ("title", "label", "claim", "text", "value"):
            value = content.get(key)
            if isinstance(value, str) and value.strip():
                return value.strip()
    component = obj.get("component_type") or obj.get("semantic_role") or obj.get("type")
    return str(component or "PPT object").replace("_", " ").title()


def _table_data(obj: dict[str, Any]) -> list[list[str]]:
    content = obj.get("content")
    if isinstance(content, dict):
        rows = content.get("rows")
        if isinstance(rows, list) and rows:
            return [[str(cell) for cell in row] for row in rows if isinstance(row, list)]
        headers = content.get("headers")
        body = content.get("body")
        if isinstance(headers, list) and isinstance(body, list):
            return [[str(cell) for cell in headers], *[[str(cell) for cell in row] for row in body if isinstance(row, list)]]
    return [["Item", "Status"], [_object_text(obj), "Editable"]]


def _chart_data(obj: dict[str, Any]) -> tuple[list[str], list[tuple[str, list[float]]]]:
    content = obj.get("content")
    if not isinstance(content, dict):
        raise ValueError(f"Chart {obj.get('id')} requires mapping content")
    categories = content.get("categories")
    raw_series = content.get("series")
    if not isinstance(categories, list) or not categories or not isinstance(raw_series, list) or not raw_series:
        raise ValueError(f"Chart {obj.get('id')} requires non-empty categories and series")
    parsed_series: list[tuple[str, list[float]]] = []
    for series in raw_series:
        if not isinstance(series, dict) or not isinstance(series.get("values"), list):
            raise ValueError(f"Chart {obj.get('id')} has invalid series")
        values = series["values"]
        if len(values) != len(categories) or not all(isinstance(value, (int, float)) for value in values):
            raise ValueError(f"Chart {obj.get('id')} series values must align with categories")
        parsed_series.append((str(series.get("name") or "Series"), [float(value) for value in values]))
    return [str(category) for category in categories], parsed_series


def _add_chart(slide: Any, obj: dict[str, Any], *, x: float, y: float, w: float, h: float, style: dict[str, Any]) -> None:
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.util import Inches, Pt

    categories, series = _chart_data(obj)
    data = ChartData()
    data.categories = categories
    for name, values in series:
        data.add_series(name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(x), Inches(y), Inches(max(w, 5.8)), Inches(h), data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.name = style["font_name"]
    chart.legend.font.size = Pt(style["table_size"])
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.tick_labels.font.name = style["font_name"]
    chart.category_axis.tick_labels.font.size = Pt(style["table_size"])


def _process_labels(obj: dict[str, Any]) -> list[str]:
    content = obj.get("content")
    nodes = content.get("nodes") if isinstance(content, dict) else None
    if not isinstance(nodes, list) or not nodes:
        raise ValueError(f"Process {obj.get('id')} requires non-empty nodes")
    labels = [str(node.get("label")) for node in nodes if isinstance(node, dict) and node.get("label") is not None]
    if len(labels) != len(nodes):
        raise ValueError(f"Process {obj.get('id')} requires a label for every node")
    return labels


def _add_process(slide: Any, obj: dict[str, Any], *, x: float, y: float, w: float, h: float, style: dict[str, Any]) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt

    labels = _process_labels(obj)
    total_width = max(w, 5.8)
    gap = 0.28
    node_width = (total_width - gap * (len(labels) - 1)) / len(labels)
    shapes = []
    for index, label in enumerate(labels):
        left = x + index * (node_width + gap)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(y), Inches(node_width), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(style["surface"])
        shape.line.color.rgb = RGBColor.from_string(style["border"])
        shape.text = label
        shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.name = style["font_name"]
            run.font.size = Pt(style["body_size"])
            run.font.color.rgb = RGBColor.from_string(style["text_primary"])
        shapes.append(shape)
    for before, after in zip(shapes, shapes[1:]):
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            before.left + before.width,
            before.top + before.height // 2,
            after.left,
            after.top + after.height // 2,
        )
        connector.line.color.rgb = RGBColor.from_string(style["primary"])


def _add_table(slide: Any, obj: dict[str, Any], *, x: float, y: float, w: float, h: float, style: dict[str, Any]) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    data = _table_data(obj)
    rows = max(len(data), 1)
    cols = max(len(row) for row in data) if data else 1
    table_shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for row_index, row in enumerate(data):
        for col_index in range(cols):
            cell = table.cell(row_index, col_index)
            cell.text = row[col_index] if col_index < len(row) else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(style["primary"] if row_index == 0 else style["background"])
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(style["table_size"])
                    run.font.name = style["font_name"]
                    run.font.color.rgb = RGBColor.from_string(style["background"] if row_index == 0 else style["text_primary"])


def _hex(value: Any, default: str) -> str:
    if isinstance(value, str) and len(value.lstrip("#")) == 6:
        return value.lstrip("#").upper()
    return default


def _resolve_style(contract: dict[str, Any]) -> dict[str, Any]:
    colors = contract.get("colors") if isinstance(contract.get("colors"), dict) else {}
    typography = contract.get("typography") if isinstance(contract.get("typography"), dict) else {}
    grid = contract.get("grid") if isinstance(contract.get("grid"), dict) else {}
    title_sizes = typography.get("title_sizes_pt") if isinstance(typography.get("title_sizes_pt"), dict) else {}
    body_sizes = typography.get("body_sizes_pt") if isinstance(typography.get("body_sizes_pt"), dict) else {}
    fonts = typography.get("font_primary") if isinstance(typography.get("font_primary"), list) else []
    return {
        "background": _hex(colors.get("background"), "FFFFFF"), "primary": _hex(colors.get("primary"), "2D3340"),
        "surface": _hex(colors.get("surface_1"), "F4F1EA"), "border": _hex(colors.get("border"), "E4E0D7"),
        "text_primary": _hex(colors.get("text_primary"), "2D3340"), "text_secondary": _hex(colors.get("text_secondary"), "6E6A60"),
        "font_name": next((font for font in fonts if isinstance(font, str) and font not in {"sans-serif", "serif", "monospace"}), "Aptos"),
        "title_size": float(title_sizes.get("slide", 22)), "body_size": float(body_sizes.get("normal", 13)),
        "table_size": float(body_sizes.get("small", 10)), "margin_left": float(grid.get("margin_left_in", 0.55)),
        "margin_top": float(grid.get("margin_top_in", 0.35)),
        "card_gap": _spacing_value(contract, "card_gap", 0.16),
    }


def _spacing_value(contract: dict[str, Any], rule: str, default: float) -> float:
    spacing = contract.get("spacing") if isinstance(contract.get("spacing"), dict) else {}
    rules = spacing.get("rules") if isinstance(spacing.get("rules"), dict) else {}
    scale = spacing.get("scale") if isinstance(spacing.get("scale"), dict) else {}
    token = rules.get(rule)
    value = scale.get(token) if isinstance(token, str) else None
    return float(value) if isinstance(value, (int, float)) else default


def _unsupported_style_warnings(contract: dict[str, Any]) -> list[str]:
    spacing = contract.get("spacing") if isinstance(contract.get("spacing"), dict) else {}
    rules = spacing.get("rules") if isinstance(spacing.get("rules"), dict) else {}
    card_gap_token = rules.get("card_gap")
    consumed_scale = {card_gap_token: True} if isinstance(card_gap_token, str) else {}
    consumed_paths: dict[str, Any] = {
        "colors": {
            "background": True,
            "primary": True,
            "surface_1": True,
            "border": True,
            "text_primary": True,
            "text_secondary": True,
        },
        "typography": {
            "font_primary": True,
            "title_sizes_pt": {"slide": True},
            "body_sizes_pt": {"normal": True, "small": True},
        },
        "grid": {"margin_left_in": True, "margin_top_in": True},
        "spacing": {
            "scale": consumed_scale,
            "rules": {"card_gap": True},
        },
    }
    warning_paths: list[str] = []
    _collect_unconsumed_style_paths(contract, consumed_paths, "", warning_paths)
    return [f"Unsupported style field ignored by PythonPptxAdapter: {path}" for path in sorted(set(warning_paths))]


def _collect_unconsumed_style_paths(value: Any, consumed_paths: Any, prefix: str, warnings: list[str]) -> None:
    if consumed_paths is True:
        return
    if isinstance(value, dict):
        if not value and prefix:
            warnings.append(prefix)
            return
        support = consumed_paths if isinstance(consumed_paths, dict) else {}
        for key, child in value.items():
            path = f"{prefix}.{key}" if prefix else key
            _collect_unconsumed_style_paths(child, support.get(key), path, warnings)
        return
    if prefix:
        warnings.append(prefix)
