from __future__ import annotations

import math
import zipfile
from pathlib import Path
from typing import Any, Optional

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .models import InspectedObject, IssueFactory, PackageInspection, SlideInspection, Thresholds

EMU_PER_INCH = 914400
R_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main", "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships"}


def emu_to_in(value: int) -> float:
    return round(float(value) / EMU_PER_INCH, 4)


def _hex(rgb: Any) -> Optional[str]:
    if rgb is None:
        return None
    text = str(rgb)
    if len(text) == 6:
        return f"#{text.upper()}"
    return None


def _shape_color(container: Any) -> list[str]:
    colors: list[str] = []
    try:
        if container is not None:
            color = _hex(container.fore_color.rgb)
            if color:
                colors.append(color)
    except Exception:
        return colors
    return colors


def _text_font_data(shape: Any) -> tuple[list[Optional[float]], list[str]]:
    sizes: list[Optional[float]] = []
    families: list[str] = []
    if not getattr(shape, "has_text_frame", False):
        return sizes, families
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            size = run.font.size
            sizes.append(round(size.pt, 2) if size is not None else None)
            if run.font.name:
                families.append(run.font.name)
    return sizes, sorted(set(families))


def _classify_shape(shape: Any, svg_shape_ids: set[int]) -> str:
    shape_id = int(getattr(shape, "shape_id", 0))
    if shape_id in svg_shape_ids:
        return "svg"
    if getattr(shape, "has_table", False):
        return "table"
    if getattr(shape, "has_chart", False):
        return "chart"
    shape_type = getattr(shape, "shape_type", None)
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        return "group"
    if shape_type in {MSO_SHAPE_TYPE.DIAGRAM, MSO_SHAPE_TYPE.IGX_GRAPHIC}:
        return "smartart"
    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if shape_type == MSO_SHAPE_TYPE.LINE:
        return "connector"
    if shape_type == MSO_SHAPE_TYPE.MEDIA:
        return "media"
    if shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
        return "ole"
    if getattr(shape, "has_text_frame", False):
        return "text"
    if shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and getattr(shape, "has_text_frame", False):
        return "text"
    return "shape"


def _collect_svg_shape_ids(pptx_path: Path, slide_index: int) -> set[int]:
    ids: set[int] = set()
    slide_part = f"ppt/slides/slide{slide_index}.xml"
    try:
        with zipfile.ZipFile(pptx_path) as package:
            root = etree.fromstring(package.read(slide_part))
    except Exception:
        return ids
    for pic in root.xpath("//p:pic", namespaces={"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}):
        nv = pic.xpath(".//p:cNvPr", namespaces={"p": "http://schemas.openxmlformats.org/presentationml/2006/main"})
        if not nv:
            continue
        ext_xml = etree.tostring(pic).decode("utf-8", errors="ignore")
        if "svgBlip" in ext_xml or ".svg" in ext_xml.lower():
            try:
                ids.add(int(nv[0].get("id")))
            except (TypeError, ValueError):
                pass
    return ids


def _image_info(shape: Any) -> tuple[Optional[str], Optional[dict[str, int]]]:
    image = getattr(shape, "image", None)
    if image is None:
        return None, None
    ext = (getattr(image, "ext", "") or "").lower()
    try:
        blob = image.blob
        from io import BytesIO

        with Image.open(BytesIO(blob)) as im:
            return ext, {"width": int(im.width), "height": int(im.height)}
    except Exception:
        return ext, None


def _iter_shapes(shapes: Any) -> list[Any]:
    collected: list[Any] = []
    for shape in shapes:
        collected.append(shape)
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            collected.extend(_iter_shapes(shape.shapes))
    return collected


def _rect_union_area(rects: list[tuple[float, float, float, float]]) -> float:
    if not rects:
        return 0.0
    xs = sorted({x1 for x1, _, x2, _ in rects} | {x2 for _, _, x2, _ in rects})
    total = 0.0
    for left, right in zip(xs, xs[1:]):
        if right <= left:
            continue
        intervals: list[tuple[float, float]] = []
        for x1, y1, x2, y2 in rects:
            if x1 < right and x2 > left:
                intervals.append((y1, y2))
        if not intervals:
            continue
        intervals.sort()
        merged: list[tuple[float, float]] = []
        for start, end in intervals:
            if not merged or start > merged[-1][1]:
                merged.append((start, end))
            else:
                merged[-1] = (merged[-1][0], max(merged[-1][1], end))
        total += (right - left) * sum(end - start for start, end in merged)
    return total


def _style_allowed_colors(style: Optional[dict[str, Any]]) -> set[str]:
    if not isinstance(style, dict):
        return set()
    colors: set[str] = set()

    def walk(value: Any) -> None:
        if isinstance(value, str) and value.startswith("#") and len(value) == 7:
            colors.add(value.upper())
        elif isinstance(value, list):
            for item in value:
                walk(item)
        elif isinstance(value, dict):
            for item in value.values():
                walk(item)

    walk(style.get("colors", {}))
    return colors


def _style_allowed_fonts(style: Optional[dict[str, Any]]) -> set[str]:
    if not isinstance(style, dict):
        return set()
    typography = style.get("typography", {})
    fonts: set[str] = set()
    for key in ("font_primary", "font_editorial", "font_mono"):
        value = typography.get(key)
        if isinstance(value, list):
            fonts.update(str(item) for item in value)
    return fonts


def _minimum_font_size(style: Optional[dict[str, Any]]) -> Optional[float]:
    if not isinstance(style, dict):
        return None
    value = style.get("typography", {}).get("minimum_body_size_pt")
    return float(value) if isinstance(value, (int, float)) else None


def _slide_expectations(ppt_ir: Optional[dict[str, Any]], slide_id: str, slide_index: int) -> tuple[int, list[str]]:
    if not isinstance(ppt_ir, dict):
        return 0, []
    for slide in ppt_ir.get("slides", []):
        if not isinstance(slide, dict):
            continue
        if slide.get("id") == slide_id or slide.get("index") == slide_index:
            texts: list[str] = []
            for field in ("title", "message"):
                if isinstance(slide.get(field), str):
                    texts.append(slide[field])
            for obj in slide.get("objects", []) or []:
                content = obj.get("content") if isinstance(obj, dict) else None
                if isinstance(content, str):
                    texts.append(content)
                elif isinstance(content, dict):
                    for value in content.values():
                        if isinstance(value, str):
                            texts.append(value)
            return sum(len(text) for text in texts), texts
    return 0, []


def _delivery_slide(delivery: Optional[dict[str, Any]], slide_id: str) -> list[dict[str, Any]]:
    if not isinstance(delivery, dict):
        return []
    for slide in delivery.get("slides", []):
        if isinstance(slide, dict) and slide.get("slide_id") == slide_id:
            return [obj for obj in slide.get("objects", []) if isinstance(obj, dict)]
    return []


def _route_allows_background(delivery_objects: list[dict[str, Any]]) -> bool:
    return any(obj.get("selected_route") == "background_image" or "background" in obj.get("rasterized_parts", []) for obj in delivery_objects)


def _route_allows_raster(delivery_objects: list[dict[str, Any]]) -> bool:
    raster_routes = {"raster_component", "generated_image", "background_image", "hybrid_overlay"}
    return any(obj.get("selected_route") in raster_routes or obj.get("rasterized_parts") for obj in delivery_objects)


def _estimate_text_overflow(obj: InspectedObject) -> Optional[float]:
    if not obj.text.strip() or obj.w_in <= 0 or obj.h_in <= 0:
        return None
    explicit = [size for size in obj.font_sizes_pt if isinstance(size, (int, float))]
    font_size = float(max(explicit)) if explicit else 14.0
    avg_char_width_in = font_size / 72.0 * 0.52
    line_height_in = font_size / 72.0 * 1.25
    chars_per_line = max(1, int(obj.w_in / max(avg_char_width_in, 0.001)))
    lines = math.ceil(len(obj.text) / chars_per_line)
    required_h = lines * line_height_in
    return required_h / obj.h_in


def inspect_slides(
    pptx_path: Path,
    package_inspection: PackageInspection,
    *,
    ppt_ir: Optional[dict[str, Any]] = None,
    style: Optional[dict[str, Any]] = None,
    delivery: Optional[dict[str, Any]] = None,
    thresholds: Optional[Thresholds] = None,
    include_raw_xml: bool = False,
) -> PackageInspection:
    thresholds = thresholds or Thresholds()
    factory = IssueFactory("pptx-structural-inspector")
    prs = Presentation(str(pptx_path))
    slide_w = emu_to_in(prs.slide_width)
    slide_h = emu_to_in(prs.slide_height)
    slide_area = max(slide_w * slide_h, 0.001)
    allowed_colors = _style_allowed_colors(style)
    allowed_fonts = _style_allowed_fonts(style)
    min_font = _minimum_font_size(style)

    deck_metrics = {
        "native_text_character_count": 0,
        "expected_text_character_count": 0,
        "native_text_ratio": None,
        "native_table_count": 0,
        "native_chart_count": 0,
        "native_connector_count": 0,
        "rasterized_area_ratio": 0.0,
        "whole_slide_raster_count": 0,
        "average_object_count": 0.0,
    }

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_id = f"S{slide_index:02d}"
        delivery_objects = _delivery_slide(delivery, slide_id)
        svg_ids = _collect_svg_shape_ids(pptx_path, slide_index)
        objects: list[InspectedObject] = []
        image_rects: list[tuple[float, float, float, float]] = []
        shapes = _iter_shapes(slide.shapes)
        for idx, shape in enumerate(shapes, start=1):
            x = emu_to_in(shape.left)
            y = emu_to_in(shape.top)
            w = emu_to_in(shape.width)
            h = emu_to_in(shape.height)
            kind = _classify_shape(shape, svg_ids)
            sizes, families = _text_font_data(shape)
            ext, image_px = _image_info(shape) if kind in {"picture", "svg"} else (None, None)
            obj = InspectedObject(
                object_id=f"ppt-object-{slide_index}-{idx}",
                slide_id=slide_id,
                shape_id=int(getattr(shape, "shape_id", idx)),
                shape_type=kind,
                name=getattr(shape, "name", "") or "",
                x_in=x,
                y_in=y,
                w_in=w,
                h_in=h,
                area_ratio=round((max(w, 0) * max(h, 0)) / slide_area, 6),
                text=getattr(shape, "text", "") if getattr(shape, "has_text_frame", False) else "",
                font_sizes_pt=sizes,
                fill_colors=_shape_color(getattr(shape, "fill", None)),
                line_colors=_shape_color(getattr(shape, "line", None)),
                source_relationship=None,
                font_families=families,
                image_ext=ext,
                image_px=image_px,
                raw={"xml": shape.element.xml} if include_raw_xml else {},
            )
            objects.append(obj)
            if kind in {"picture", "svg"}:
                image_rects.append((max(0.0, x), max(0.0, y), min(slide_w, x + w), min(slide_h, y + h)))

        expected_chars, expected_texts = _slide_expectations(ppt_ir, slide_id, slide_index)
        native_chars = sum(len(obj.text.strip()) for obj in objects if obj.shape_type == "text")
        text_boxes = [obj for obj in objects if obj.shape_type == "text"]
        tiny = [
            obj
            for obj in objects
            if obj.w_in < thresholds.tiny_width_or_height_in
            or obj.h_in < thresholds.tiny_width_or_height_in
            or (obj.w_in * obj.h_in) < thresholds.tiny_area_sq_in
        ]
        raster_area = _rect_union_area(image_rects)
        metrics = {
            "native_text_character_count": native_chars,
            "expected_text_character_count": expected_chars,
            "native_text_ratio": round(min(native_chars, expected_chars) / expected_chars, 4) if expected_chars else None,
            "text_box_count": len(text_boxes),
            "average_characters_per_text_box": round(native_chars / len(text_boxes), 4) if text_boxes else 0,
            "small_text_box_count": len([obj for obj in text_boxes if obj.w_in < 1.0 or obj.h_in < 0.25]),
            "font_family_count": len({family for obj in text_boxes for family in obj.font_families}),
            "font_size_count": len({size for obj in text_boxes for size in obj.font_sizes_pt if size is not None}),
            "object_count": len(objects),
            "tiny_object_count": len(tiny),
            "text_boxes_per_100_characters": round(len(text_boxes) / max(native_chars, 1) * 100, 4),
            "rasterized_area_ratio": round(raster_area / slide_area, 6),
            "native_table_count": len([obj for obj in objects if obj.shape_type == "table"]),
            "native_chart_count": len([obj for obj in objects if obj.shape_type == "chart"]),
            "native_connector_count": len([obj for obj in objects if obj.shape_type == "connector"]),
        }
        slide_result = SlideInspection(slide_index=slide_index, slide_id=slide_id, objects=objects, metrics=metrics)
        _add_slide_issues(slide_result, factory, slide_w, slide_h, thresholds, allowed_colors, allowed_fonts, min_font, delivery_objects, expected_texts)
        package_inspection.slides.append(slide_result)

        deck_metrics["native_text_character_count"] += native_chars
        deck_metrics["expected_text_character_count"] += expected_chars
        deck_metrics["native_table_count"] += metrics["native_table_count"]
        deck_metrics["native_chart_count"] += metrics["native_chart_count"]
        deck_metrics["native_connector_count"] += metrics["native_connector_count"]
        deck_metrics["rasterized_area_ratio"] += raster_area / slide_area
        deck_metrics["whole_slide_raster_count"] += len([issue for issue in slide_result.issues if issue.issue_code == "PPTX_WHOLE_SLIDE_RASTER"])

    slide_count = max(len(prs.slides), 1)
    deck_metrics["native_text_ratio"] = (
        round(min(deck_metrics["native_text_character_count"], deck_metrics["expected_text_character_count"]) / deck_metrics["expected_text_character_count"], 4)
        if deck_metrics["expected_text_character_count"]
        else None
    )
    deck_metrics["rasterized_area_ratio"] = round(deck_metrics["rasterized_area_ratio"] / slide_count, 6)
    deck_metrics["average_object_count"] = round(sum(slide.metrics["object_count"] for slide in package_inspection.slides) / slide_count, 4)
    package_inspection.metrics.update(deck_metrics)
    package_inspection.status = "failed" if any(issue.severity in {"error", "fatal"} for issue in _all_issues(package_inspection)) else "passed"
    return package_inspection


def _all_issues(inspection: PackageInspection) -> list[Any]:
    issues: list[Any] = list(inspection.issues)
    for slide in inspection.slides:
        issues.extend(slide.issues)
    return issues


def _add_slide_issues(
    slide: SlideInspection,
    factory: IssueFactory,
    slide_w: float,
    slide_h: float,
    thresholds: Thresholds,
    allowed_colors: set[str],
    allowed_fonts: set[str],
    min_font: Optional[float],
    delivery_objects: list[dict[str, Any]],
    expected_texts: list[str],
) -> None:
    if not slide.objects:
        slide.issues.append(factory.create("PPTX_BLANK_SLIDE", "error", "content", "Slide contains no editable or visual objects.", slide_id=slide.slide_id, evidence={}))
    for obj in slide.objects:
        declared_non_text_layer = obj.name.lower().startswith(("decoration:", "background:", "material:", "connector:"))
        if (
            obj.shape_type == "text"
            and not obj.text.strip()
            and obj.fill_colors
            and obj.area_ratio >= 0.01
            and not declared_non_text_layer
        ):
            slide.issues.append(factory.create(
                "PPTX_ORPHAN_SOLID_COLOR_BLOCK",
                "error",
                "visual",
                "Unlabeled solid-color block may be a stray arrowhead or accidental shape.",
                slide_id=slide.slide_id,
                object_id=obj.object_id,
                ppt_shape_id=obj.shape_id,
                evidence={"name": obj.name, "fill_colors": obj.fill_colors, "area_ratio": obj.area_ratio},
            ))
        if obj.shape_type == "text" and not obj.text.strip() and not declared_non_text_layer:
            slide.issues.append(factory.create("PPTX_EMPTY_TEXT_BOX", "warning", "editability", "Text box is empty.", slide_id=slide.slide_id, object_id=obj.object_id, ppt_shape_id=obj.shape_id, evidence={"name": obj.name}))
        if obj.x_in < 0 or obj.y_in < 0 or obj.x_in + obj.w_in > slide_w or obj.y_in + obj.h_in > slide_h:
            slide.issues.append(factory.create("PPTX_TEXT_OUT_OF_BOUNDS" if obj.shape_type == "text" else "GEOMETRY_OBJECT_OUT_OF_BOUNDS", "error", "geometry", "Object extends beyond slide bounds.", slide_id=slide.slide_id, object_id=obj.object_id, ppt_shape_id=obj.shape_id, evidence={"x_in": obj.x_in, "y_in": obj.y_in, "w_in": obj.w_in, "h_in": obj.h_in, "slide_w_in": slide_w, "slide_h_in": slide_h}))
        if min_font is not None:
            for size in obj.font_sizes_pt:
                if isinstance(size, (int, float)) and size < min_font:
                    slide.issues.append(factory.create("PPTX_FONT_SIZE_BELOW_MINIMUM", "error", "style", "Font size is below the style contract minimum.", slide_id=slide.slide_id, object_id=obj.object_id, ppt_shape_id=obj.shape_id, evidence={"font_size_pt": size, "minimum_body_size_pt": min_font}))
        if allowed_fonts:
            for family in obj.font_families:
                if family not in allowed_fonts:
                    slide.issues.append(factory.create("PPTX_FONT_NOT_IN_CONTRACT", "warning", "style", "Font family is not in the style contract font stacks.", slide_id=slide.slide_id, object_id=obj.object_id, ppt_shape_id=obj.shape_id, evidence={"font_family": family, "allowed_fonts": sorted(allowed_fonts)}))
        if allowed_colors:
            for color in obj.fill_colors + obj.line_colors:
                if color.upper() not in allowed_colors:
                    slide.issues.append(factory.create("STYLE_COLOR_DRIFT", "warning", "style", "Object uses a color outside the style contract palette.", slide_id=slide.slide_id, object_id=obj.object_id, ppt_shape_id=obj.shape_id, evidence={"color": color, "allowed_colors": sorted(allowed_colors)}))
        overflow = _estimate_text_overflow(obj)
        if overflow is not None and overflow > 1.15:
            slide.issues.append(factory.create("TEXT_OVERFLOW_RISK", "warning", "text", "Text box may overflow based on character count, font size, and box geometry.", slide_id=slide.slide_id, object_id=obj.object_id, ppt_shape_id=obj.shape_id, evidence={"overflow_ratio_estimate": round(overflow, 4)}))
        if obj.shape_type in {"picture", "svg"} and obj.image_px:
            ppi_x = obj.image_px["width"] / max(obj.w_in, 0.001)
            ppi_y = obj.image_px["height"] / max(obj.h_in, 0.001)
            if min(ppi_x, ppi_y) < thresholds.low_resolution_ppi:
                slide.issues.append(factory.create("PPTX_LOW_RESOLUTION_IMAGE", "warning", "editability", "Image resolution is low for its displayed size.", slide_id=slide.slide_id, object_id=obj.object_id, ppt_shape_id=obj.shape_id, evidence={"ppi_x": round(ppi_x, 2), "ppi_y": round(ppi_y, 2), "image_px": obj.image_px}))

    object_count = slide.metrics["object_count"]
    if object_count > thresholds.objects_per_slide_error:
        slide.issues.append(factory.create("PPTX_EXCESSIVE_OBJECT_COUNT", "error", "editability", "Slide has too many objects.", slide_id=slide.slide_id, evidence={"object_count": object_count, "threshold": thresholds.objects_per_slide_error}))
    elif object_count > thresholds.objects_per_slide_warn:
        slide.issues.append(factory.create("PPTX_OBJECT_FRAGMENTATION", "warning", "editability", "Slide object count is high.", slide_id=slide.slide_id, evidence={"object_count": object_count, "threshold": thresholds.objects_per_slide_warn}))
    tiny_count = slide.metrics["tiny_object_count"]
    if tiny_count > thresholds.tiny_object_error:
        slide.issues.append(factory.create("PPTX_TINY_OBJECT_OVERLOAD", "error", "editability", "Slide has too many tiny objects.", slide_id=slide.slide_id, evidence={"tiny_object_count": tiny_count, "threshold": thresholds.tiny_object_error}))
    elif tiny_count > thresholds.tiny_object_warn:
        slide.issues.append(factory.create("PPTX_TINY_OBJECT_OVERLOAD", "warning", "editability", "Slide has many tiny objects.", slide_id=slide.slide_id, evidence={"tiny_object_count": tiny_count, "threshold": thresholds.tiny_object_warn}))
    fragmentation = slide.metrics["text_boxes_per_100_characters"]
    if fragmentation > thresholds.text_boxes_per_100_characters_error:
        slide.issues.append(factory.create("PPTX_TEXT_FRAGMENTATION", "error", "editability", "Text is split across too many text boxes.", slide_id=slide.slide_id, evidence={"text_boxes_per_100_characters": fragmentation}))
    elif fragmentation > thresholds.text_boxes_per_100_characters_warn:
        slide.issues.append(factory.create("PPTX_TEXT_FRAGMENTATION", "warning", "editability", "Text box fragmentation is high.", slide_id=slide.slide_id, evidence={"text_boxes_per_100_characters": fragmentation}))

    native_text = "\n".join(obj.text for obj in slide.objects if obj.shape_type == "text")
    for expected in expected_texts:
        if expected and expected not in native_text:
            slide.issues.append(factory.create("PPTX_EXPECTED_TEXT_MISSING", "error", "editability", "Expected PPT IR text is missing from native text objects.", slide_id=slide.slide_id, evidence={"expected_text": expected[:160]}))

    pictures = [obj for obj in slide.objects if obj.shape_type in {"picture", "svg"}]
    text_count = slide.metrics["native_text_character_count"]
    for obj in pictures:
        covers_center = obj.x_in <= slide_w / 2 <= obj.x_in + obj.w_in and obj.y_in <= slide_h / 2 <= obj.y_in + obj.h_in
        is_background = "background" in obj.name.lower()
        background_allowed = is_background and text_count > 20 and _route_allows_background(delivery_objects)
        if obj.area_ratio >= thresholds.whole_slide_raster_area_ratio and covers_center and not background_allowed:
            slide.issues.append(factory.create("PPTX_WHOLE_SLIDE_RASTER", "error", "editability", "Detected a near full-slide raster image without native overlay exemption.", slide_id=slide.slide_id, object_id=obj.object_id, ppt_shape_id=obj.shape_id, evidence={"image_area_ratio": obj.area_ratio, "native_text_character_count": text_count}))
        elif obj.area_ratio >= thresholds.large_image_area_ratio and not _route_allows_raster(delivery_objects):
            slide.issues.append(factory.create("PPTX_LARGE_UNDECLARED_IMAGE", "warning", "delivery", "Large image is not declared by the delivery plan.", slide_id=slide.slide_id, object_id=obj.object_id, ppt_shape_id=obj.shape_id, evidence={"image_area_ratio": obj.area_ratio}))

    planned_routes = [obj.get("selected_route") for obj in delivery_objects]
    if "native_table" in planned_routes and slide.metrics["native_table_count"] == 0:
        slide.issues.append(factory.create("PPTX_TABLE_NOT_NATIVE", "error", "delivery", "Delivery plan selected native_table but no native table was found.", slide_id=slide.slide_id, evidence={"planned_routes": planned_routes}))
    if "native_chart" in planned_routes and slide.metrics["native_chart_count"] == 0:
        slide.issues.append(factory.create("PPTX_CHART_NOT_NATIVE", "error", "delivery", "Delivery plan selected native_chart but no native chart was found.", slide_id=slide.slide_id, evidence={"planned_routes": planned_routes}))
    if pictures and any(route in {"native_table", "native_chart", "native_diagram"} for route in planned_routes):
        slide.issues.append(factory.create("PPTX_RASTER_ROUTE_UNDECLARED", "error", "delivery", "Raster object appears where delivery plan selected a native route.", slide_id=slide.slide_id, evidence={"planned_routes": planned_routes, "picture_count": len(pictures)}))
