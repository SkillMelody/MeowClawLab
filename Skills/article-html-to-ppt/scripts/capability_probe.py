#!/usr/bin/env python3
"""Probe local PPT builder, renderer, format, and font capabilities."""

from __future__ import annotations

import argparse
import importlib.metadata
import json
import os
import platform
import shutil
import subprocess
import sys
import tempfile
import xml.etree.ElementTree as ET
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from ppt_qa.renderers.libreoffice import LibreOfficeRenderer
from ppt_qa.renderers.powerpoint_macos import PowerPointMacOSRenderer

FEATURE_DEFAULTS: dict[str, bool | str] = {
    "native_text": "unknown",
    "native_table": "unknown",
    "native_chart": "unknown",
    "native_connector": "unknown",
    "svg_embed": "unknown",
    "speaker_notes": "unknown",
    "theme": "unknown",
    "slide_master": "unknown",
    "readback": "unknown",
}
GENERIC_FONTS = {"sans-serif", "serif", "monospace", "cursive", "fantasy", "system-ui"}


def load_json(path: Path | None) -> Any:
    if not path:
        return None
    with path.open("r", encoding="utf-8") as handle:
        return json.load(handle)


def write_json(path: Path, data: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def run_command(command: list[str], timeout: int) -> subprocess.CompletedProcess[str]:
    return subprocess.run(command, text=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=timeout, check=False)


def version_from_command(command: list[str], timeout: int) -> str | None:
    try:
        completed = run_command(command, timeout)
    except (OSError, subprocess.TimeoutExpired):
        return None
    output = (completed.stdout or completed.stderr).strip()
    return output.splitlines()[0] if output else None


def command_path(name: str) -> str | None:
    path = shutil.which(name)
    return name if path else None


def package_version(name: str) -> str | None:
    try:
        return importlib.metadata.version(name)
    except importlib.metadata.PackageNotFoundError:
        return None


def registry_component_support(registry: dict[str, Any] | None, builder: str) -> dict[str, str]:
    support: dict[str, str] = {}
    for component in (registry or {}).get("components", []) or []:
        if not isinstance(component, dict):
            continue
        levels = component.get("builder_support", {})
        if builder in levels:
            support[component["component_type"]] = str(levels[builder].get("level", "unknown"))
    return support


def probe_python_pptx(registry: dict[str, Any] | None, timeout: int) -> dict[str, Any]:
    warnings: list[str] = []
    errors: list[str] = []
    version = package_version("python-pptx")
    available = False
    if version:
        try:
            with tempfile.TemporaryDirectory() as tmp:
                script = (
                    "from pptx import Presentation\n"
                    "p=Presentation()\n"
                    "p.slides.add_slide(p.slide_layouts[6])\n"
                    f"p.save({str(Path(tmp) / 'smoke.pptx')!r})\n"
                )
                completed = run_command([sys.executable, "-c", script], timeout)
                available = completed.returncode == 0 and (Path(tmp) / "smoke.pptx").exists()
                if not available:
                    errors.append("CAPABILITY_SMOKE_TEST_FAILED")
        except (OSError, subprocess.TimeoutExpired):
            errors.append("CAPABILITY_SMOKE_TEST_FAILED")
    else:
        warnings.append("python-pptx package is not importable.")
    return {
        "available": available,
        "version": version,
        "command": sys.executable if version else None,
        "features": {
            "native_text": available,
            "native_table": available,
            "native_chart": False,
            "native_connector": "partial" if available else False,
            "svg_embed": False,
            "speaker_notes": False,
            "theme": "partial" if available else False,
            "slide_master": "partial" if available else False,
            "readback": available,
        },
        "components": registry_component_support(registry, "python_pptx") if available else {},
        "warnings": warnings,
        "errors": errors,
    }


def probe_pptxgenjs(registry: dict[str, Any] | None, timeout: int) -> dict[str, Any]:
    node = shutil.which("node")
    warnings: list[str] = []
    errors: list[str] = []
    version = None
    available = False
    if not node:
        warnings.append("node command not found.")
    else:
        script = (
            "const pptxgen = require('pptxgenjs');"
            "const pkg = require('pptxgenjs/package.json');"
            "const pptx = new pptxgen();"
            "pptx.layout = 'LAYOUT_WIDE';"
            "let slide = pptx.addSlide();"
            "slide.addText('smoke', {x:0.2,y:0.2,w:1,h:0.3,fontSize:10});"
            "console.log(pkg.version || 'unknown');"
        )
        try:
            completed = run_command([node, "-e", script], timeout)
            if completed.returncode == 0:
                available = True
                version = (completed.stdout or "").strip().splitlines()[-1] if completed.stdout.strip() else None
            else:
                warnings.append("pptxgenjs package is not importable by node.")
                errors.append("CAPABILITY_SMOKE_TEST_FAILED")
        except (OSError, subprocess.TimeoutExpired):
            errors.append("CAPABILITY_SMOKE_TEST_FAILED")
    return {
        "available": available,
        "version": version,
        "command": "node" if node else None,
        "features": {
            "native_text": available,
            "native_table": available,
            "native_chart": available,
            "native_connector": available,
            "svg_embed": available,
            "speaker_notes": available,
            "theme": "partial" if available else False,
            "slide_master": "partial" if available else False,
            "readback": False,
        },
        "components": registry_component_support(registry, "pptxgenjs") if available else {},
        "warnings": warnings,
        "errors": errors,
    }


def probe_officecli(registry: dict[str, Any] | None, timeout: int) -> dict[str, Any]:
    officecli = shutil.which("officecli")
    warnings: list[str] = []
    errors: list[str] = []
    version = None
    available = False
    if not officecli:
        warnings.append("officecli command not found.")
    else:
        version = version_from_command([officecli, "--version"], timeout)
        try:
            completed = run_command([officecli, "--help"], timeout)
            available = completed.returncode == 0 and bool(version or completed.stdout)
            if not available:
                errors.append("CAPABILITY_SMOKE_TEST_FAILED")
        except (OSError, subprocess.TimeoutExpired):
            errors.append("CAPABILITY_SMOKE_TEST_FAILED")
    return {
        "available": available,
        "version": version,
        "command": "officecli" if officecli else None,
        "features": FEATURE_DEFAULTS | ({"readback": available} if available else {}),
        "components": registry_component_support(registry, "officecli") if available else {},
        "warnings": warnings,
        "errors": errors,
    }


def probe_html_svg(registry: dict[str, Any] | None, timeout: int) -> dict[str, Any]:
    del timeout
    warnings: list[str] = []
    errors: list[str] = []
    available = False
    try:
        svg = "<svg xmlns='http://www.w3.org/2000/svg' width='16' height='9'><text x='1' y='8'>smoke</text></svg>"
        ET.fromstring(svg)
        available = True
    except ET.ParseError:
        errors.append("CAPABILITY_SMOKE_TEST_FAILED")
    return {
        "available": available,
        "version": "python-stdlib-svg" if available else None,
        "command": "python",
        "features": {
            "native_text": False,
            "native_table": False,
            "native_chart": False,
            "native_connector": False,
            "svg_embed": available,
            "speaker_notes": False,
            "theme": False,
            "slide_master": False,
            "readback": False,
        },
        "components": registry_component_support(registry, "html_svg") if available else {},
        "warnings": warnings + ["html_svg is visual-only; it cannot satisfy native-required objects."],
        "errors": errors,
    }


def probe_builders(registry: dict[str, Any] | None, timeout: int, only: str | None = None) -> dict[str, Any]:
    probes = {
        "officecli": probe_officecli,
        "pptxgenjs": probe_pptxgenjs,
        "python_pptx": probe_python_pptx,
        "html_svg": probe_html_svg,
    }
    selected = [only] if only else list(probes)
    return {name: probes[name](registry, timeout) for name in selected if name in probes}


def probe_renderers(only: str | None = None) -> dict[str, Any]:
    power = PowerPointMacOSRenderer()
    libre = LibreOfficeRenderer()
    renderers = {
        "microsoft_powerpoint": {
            "available": power.is_available(),
            "version": power.version() if power.is_available() else None,
            "command": "Microsoft PowerPoint.app" if power.is_available() else None,
            "supports_pdf": power.is_available(),
            "supports_png": power.is_available(),
            "supports_readback": power.is_available(),
        },
        "libreoffice": {
            "available": libre.is_available(),
            "version": libre.version() if libre.is_available() else None,
            "command": "soffice" if libre.is_available() else None,
            "supports_pdf": libre.is_available(),
            "supports_png": libre.is_available(),
            "supports_readback": False,
        },
        "keynote": {
            "available": Path("/Applications/Keynote.app").exists(),
            "version": None,
            "command": "Keynote.app" if Path("/Applications/Keynote.app").exists() else None,
            "supports_pdf": False,
            "supports_png": False,
            "supports_readback": False,
        },
    }
    if only:
        return {only: renderers[only]} if only in renderers else {}
    return renderers


def list_fonts(timeout: int) -> list[str]:
    names: set[str] = set()
    system = platform.system().lower()
    commands: list[list[str]] = []
    if shutil.which("fc-list"):
        commands.append(["fc-list", ":", "family"])
    if system == "darwin":
        commands.append(["system_profiler", "SPFontsDataType"])
    for command in commands:
        try:
            completed = run_command(command, timeout)
        except (OSError, subprocess.TimeoutExpired):
            continue
        if completed.returncode != 0:
            continue
        for line in completed.stdout.splitlines():
            text = line.strip()
            if not text:
                continue
            if command[0] == "fc-list":
                for family in text.split(","):
                    if family.strip():
                        names.add(family.strip())
            elif text.endswith(":") and not text.startswith(" "):
                names.add(text[:-1].strip())
            elif "Full Name:" in text:
                names.add(text.split("Full Name:", 1)[1].strip())
    return sorted(names)


def required_fonts(style: dict[str, Any] | None) -> list[str]:
    fonts: list[str] = []
    typography = (style or {}).get("typography", {}) if isinstance(style, dict) else {}
    for key, stack in typography.items():
        if key.startswith("font_") and isinstance(stack, list):
            fonts.extend(str(item) for item in stack if isinstance(item, str) and item not in GENERIC_FONTS)
    return list(dict.fromkeys(fonts))


def font_stacks(style: dict[str, Any] | None) -> list[list[str]]:
    stacks: list[list[str]] = []
    typography = (style or {}).get("typography", {}) if isinstance(style, dict) else {}
    for key, stack in typography.items():
        if key.startswith("font_") and isinstance(stack, list):
            stacks.append([str(item) for item in stack if isinstance(item, str)])
    return stacks


def probe_fonts(style: dict[str, Any] | None, timeout: int) -> dict[str, Any]:
    installed = list_fonts(timeout)
    installed_lower = {font.lower() for font in installed}
    missing = [font for font in required_fonts(style) if font.lower() not in installed_lower]
    fallback_map: dict[str, str] = {}
    for stack in font_stacks(style):
        for idx, font in enumerate(stack):
            if font not in missing:
                continue
            later = stack[idx + 1 :]
            fallback = next((candidate for candidate in later if candidate.lower() in installed_lower), None)
            if fallback is None:
                fallback = next((candidate for candidate in later if candidate in GENERIC_FONTS), None)
            if fallback:
                fallback_map[font] = fallback
    return {
        "installed": installed[:250],
        "required_missing": missing,
        "fallback_map": fallback_map,
    }


def importable(name: str) -> bool:
    try:
        __import__(name)
        return True
    except Exception:
        return False


def probe_formats(builders: dict[str, Any], renderers: dict[str, Any]) -> dict[str, Any]:
    svg_renderer = bool(shutil.which("rsvg-convert") or shutil.which("inkscape") or shutil.which("magick") or importable("cairosvg"))
    pdf_to_png = bool(shutil.which("pdftoppm") or shutil.which("magick") or importable("fitz"))
    pptx_python = bool(builders.get("python_pptx", {}).get("available"))
    pptxgen = bool(builders.get("pptxgenjs", {}).get("available"))
    renderer_available = any(renderer.get("available") for renderer in renderers.values())
    return {
        "html_render": bool(shutil.which("chromium") or shutil.which("google-chrome") or shutil.which("playwright") or importable("playwright")),
        "html_to_svg": bool(builders.get("html_svg", {}).get("available")),
        "svg_render": svg_renderer,
        "svg_embed": pptxgen or pptx_python,
        "pdf_to_png": pdf_to_png,
        "pptx_read": pptx_python,
        "pptx_write": pptx_python or pptxgen,
        "pptx_render": renderer_available,
        "image_generation": False,
    }


def environment() -> dict[str, Any]:
    return {
        "os": platform.system() or None,
        "os_version": platform.release() or None,
        "architecture": platform.machine() or None,
        "python": platform.python_version(),
        "node": version_from_command(["node", "--version"], 3) if shutil.which("node") else None,
        "java": version_from_command(["java", "-version"], 3) if shutil.which("java") else None,
        "working_directory": Path.cwd().name,
        "containerized": Path("/.dockerenv").exists() or os.environ.get("container") is not None,
    }


def collect_messages(report: dict[str, Any]) -> None:
    warnings: list[str] = []
    errors: list[str] = []
    for name, builder in report["builders"].items():
        for warning in builder.get("warnings", []) or []:
            warnings.append(f"{name}: {warning}")
        for error in builder.get("errors", []) or []:
            if builder.get("available"):
                errors.append(f"{name}: {error}")
            else:
                warnings.append(f"{name}: {error}")
    if not any(renderer.get("available") for renderer in report["renderers"].values()):
        warnings.append("CAPABILITY_RENDERER_NOT_FOUND: no PowerPoint/LibreOffice renderer detected.")
    if report["fonts"].get("required_missing"):
        warnings.append("CAPABILITY_REQUIRED_FONT_MISSING: one or more style fonts are missing.")
    report["warnings"] = sorted(set(report.get("warnings", []) + warnings))
    report["errors"] = sorted(set(report.get("errors", []) + errors))


def parse_args(argv: list[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate a local PPT capability report from real environment checks.")
    parser.add_argument("--style", type=Path)
    parser.add_argument("--registry", type=Path)
    parser.add_argument("--output", type=Path, default=Path(".ppt-work/capability-report.json"))
    parser.add_argument("--json", action="store_true")
    parser.add_argument("--check-builder")
    parser.add_argument("--check-renderer")
    parser.add_argument("--check-fonts", action="store_true")
    parser.add_argument("--timeout", type=int, default=8)
    parser.add_argument("--strict", action="store_true")
    return parser.parse_args(argv)


def main(argv: list[str]) -> int:
    args = parse_args(argv)
    if args.timeout <= 0:
        print("timeout must be positive", file=sys.stderr)
        return 2
    style = load_json(args.style)
    registry = load_json(args.registry)
    builders = probe_builders(registry, args.timeout, only=args.check_builder)
    renderers = probe_renderers(only=args.check_renderer)
    fonts = probe_fonts(style, args.timeout) if (args.check_fonts or args.style) else {"installed": [], "required_missing": [], "fallback_map": {}}
    report = {
        "schema_version": "1.0",
        "generated_at": datetime.now(timezone.utc).astimezone().isoformat(timespec="seconds"),
        "environment": environment(),
        "builders": builders,
        "renderers": renderers,
        "formats": {},
        "fonts": fonts,
        "features": {},
        "warnings": [],
        "errors": [],
    }
    report["formats"] = probe_formats(builders, renderers)
    report["features"] = {
        "available_builder_count": sum(1 for builder in builders.values() if builder.get("available")),
        "available_renderer_count": sum(1 for renderer in renderers.values() if renderer.get("available")),
        "native_builder_available": any(
            builder.get("available") and builder.get("features", {}).get("native_text") is True
            for builder in builders.values()
        ),
    }
    collect_messages(report)
    write_json(args.output, report)
    if args.json:
        print(json.dumps(report, ensure_ascii=False, indent=2))
    else:
        print(f"capability report written: {args.output}")
    return 1 if args.strict and report["errors"] else 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
