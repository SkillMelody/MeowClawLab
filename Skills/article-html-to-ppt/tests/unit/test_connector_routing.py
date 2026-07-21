from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.util import Inches

from builders.python_pptx_adapter import _add_routed_connector
from ppt_qa.package_inspector import inspect_package
from ppt_qa.slide_inspector import inspect_slides


def _node(slide, x: float, y: float, w: float = 1.4, h: float = 0.7):
    return slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )


def _arrow_xml(shape) -> str:
    return etree.tostring(shape.element).decode("utf-8")


def test_skill_keeps_206_brand_name_and_advances_only_version() -> None:
    skill = (Path(__file__).resolve().parents[2] / "SKILL.md").read_text(encoding="utf-8")
    assert "display_name: MeowClaw PPT Smith" in skill
    assert "english_alias: MeowClaw PPT Smith" in skill
    assert "# MeowClaw PPT Smith" in skill
    assert "version: 2.0.7" in skill


def test_routed_connector_emits_native_straight_elbow_and_curve_with_binding() -> None:
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    a = _node(slide, 0.8, 1.0)
    b = _node(slide, 3.0, 1.0)
    c = _node(slide, 3.0, 3.0)

    straight = _add_routed_connector(slide, a, b, route="straight")
    elbow = _add_routed_connector(slide, b, c, route="orthogonal")
    curve = _add_routed_connector(slide, c, a, route="curved", feedback=True)

    assert 'prst="line"' in _arrow_xml(straight)
    assert 'prst="bentConnector' in _arrow_xml(elbow)
    assert 'prst="curvedConnector' in _arrow_xml(curve)
    for connector in (straight, elbow, curve):
        xml = _arrow_xml(connector)
        assert "stCxn" in xml
        assert "endCxn" in xml
        assert "tailEnd" in xml and 'type="triangle"' in xml


def test_inspector_rejects_connector_that_passes_through_unrelated_node(tmp_path: Path) -> None:
    path = tmp_path / "through-node.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    _node(slide, 3.0, 2.0, 2.0, 1.0)
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(1.0), Inches(2.5), Inches(7.0), Inches(2.5),
    )
    connector.name = "Connector:primary"
    tail = etree.Element(qn("a:tailEnd")); tail.set("type", "triangle")
    connector.line._get_or_add_ln().append(tail)
    deck.save(path)

    result = inspect_slides(path, inspect_package(path), include_raw_xml=True)
    codes = {issue.issue_code for slide_result in result.slides for issue in slide_result.issues}
    assert "PPTX_CONNECTOR_THROUGH_NODE" in codes


def test_inspector_rejects_connector_crossing(tmp_path: Path) -> None:
    path = tmp_path / "crossing.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(6), Inches(5)).name = "Connector:a"
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(5), Inches(6), Inches(1)).name = "Connector:b"
    deck.save(path)

    result = inspect_slides(path, inspect_package(path), include_raw_xml=True)
    codes = {issue.issue_code for slide_result in result.slides for issue in slide_result.issues}
    assert "PPTX_CONNECTOR_CROSSING" in codes
