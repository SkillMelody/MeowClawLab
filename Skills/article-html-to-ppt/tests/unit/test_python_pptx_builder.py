from __future__ import annotations

import json
import subprocess
import sys
from pathlib import Path

from jsonschema import Draft202012Validator

ROOT = Path(__file__).resolve().parents[2]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from builders.python_pptx_adapter import PythonPptxAdapter
from ppt_qa.package_inspector import inspect_package

STYLE = ROOT / "tests" / "fixtures" / "styles" / "consulting-light.json"
EDITORIAL_STYLE = ROOT / "tests" / "fixtures" / "v2-acceptance" / "style-contract-editorial.json"
BUILD = ROOT / "scripts" / "build_deck.py"


def sample_ppt_ir() -> dict:
    return {
        "schema_version": "2.0",
        "deck": {
            "id": "builder-smoke",
            "title": "Builder Smoke",
            "source_type": "article",
            "audience": "builders",
            "purpose": "smoke test",
            "language": "en",
            "presentation_context": "internal",
            "aspect_ratio": "16:9",
            "production_profile": "fast",
            "target_builder": "python_pptx",
            "logical_slide_count": 2,
        },
        "sources": [{"source_id": "src-001", "type": "article", "title": "Fixture"}],
        "style_contract_ref": "style-contract.json",
        "asset_manifest_ref": "asset-manifest.json",
        "slides": [
            {
                "id": "S01",
                "index": 1,
                "slide_role": "cover",
                "title_role": "navigation",
                "title": "Builder Smoke",
                "judgment": None,
                "message": "Native PPTX smoke test.",
                "audience_question": "Can the adapter create a deck?",
                "source_refs": [],
                "primary_expression": "textual_argument",
                "supporting_expressions": [],
                "primary_anchor": "title",
                "objects": [],
                "delivery_contract": {"preferred_route": "native_ppt", "editable_core": ["title"], "raster_allowance": [], "forbidden_raster": ["title"]},
            },
            {
                "id": "S02",
                "index": 2,
                "slide_role": "judgment",
                "title_role": "judgment",
                "title": "Text and tables are native",
                "judgment": "The minimal builder creates editable text and table objects.",
                "message": "Use it as the first v2 runtime gate.",
                "audience_question": "What does the minimal builder cover?",
                "source_refs": [{"source_id": "src-001", "locator": "s1", "claim_type": "direct"}],
                "primary_expression": "table_matrix",
                "supporting_expressions": [],
                "primary_anchor": "table",
                "objects": [
                    {
                        "id": "tbl-1",
                        "type": "table",
                        "component_type": "native_table",
                        "semantic_role": "evidence",
                        "content": {"headers": ["Capability", "Status"], "body": [["Text", "Native"], ["Table", "Native"]]},
                        "source_refs": [],
                        "editability": "native_required",
                    }
                ],
                "delivery_contract": {"preferred_route": "native_table", "editable_core": ["headers", "cells"], "raster_allowance": [], "forbidden_raster": ["tbl-1"]},
            },
        ],
    }


def sample_delivery_plan() -> dict:
    return {
        "schema_version": "1.0",
        "ppt_ir_ref": "ppt-ir.json",
        "style_contract_ref": "style-contract.json",
        "component_registry_ref": "references/component-registry.json",
        "capability_report_ref": None,
        "profile": "fast",
        "builder": {"requested": "python_pptx", "selected": "python_pptx", "version": None, "selection_score": 1, "selection_reasons": ["TEST"]},
        "slides": [
            {
                "slide_id": "S02",
                "objects": [
                    {
                        "slide_id": "S02",
                        "object_id": "tbl-1",
                        "component_type": "native_table",
                        "preferred_route": "native_table",
                        "selected_route": "native_table",
                        "decision": "selected",
                        "reason_codes": [],
                        "editable_core": ["headers", "cells"],
                        "rasterized_parts": [],
                        "svg_parts": [],
                        "native_overlay_parts": [],
                        "qa_checks": ["TABLE_NATIVE"],
                        "fallback_chain": [],
                    }
                ],
            }
        ],
        "summary": {"total_objects": 1, "fallback_count": 0, "unsupported_count": 0, "risk_codes": []},
    }


def write_json(path: Path, data: dict) -> None:
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def test_python_pptx_adapter_builds_inspectable_deck(tmp_path: Path) -> None:
    adapter = PythonPptxAdapter()
    plan = adapter.plan(sample_ppt_ir(), json.loads(STYLE.read_text(encoding="utf-8")), sample_delivery_plan())
    result = adapter.build(plan, tmp_path)
    assert result.status == "created"
    assert result.pptx is not None
    inspection = inspect_package(Path(result.pptx), ppt_ir=sample_ppt_ir())
    assert inspection.status == "passed"
    assert inspection.slide_count == 2


def test_build_deck_cli_writes_manifest(tmp_path: Path) -> None:
    ppt_ir = tmp_path / "ppt-ir.json"
    style = tmp_path / "style-contract.json"
    delivery = tmp_path / "delivery-plan.json"
    manifest = tmp_path / "build-manifest.json"
    output = tmp_path / "out"
    write_json(ppt_ir, sample_ppt_ir())
    style.write_text(STYLE.read_text(encoding="utf-8"), encoding="utf-8")
    write_json(delivery, sample_delivery_plan())

    result = subprocess.run(
        [
            sys.executable,
            str(BUILD),
            "--ppt-ir",
            str(ppt_ir),
            "--style",
            str(style),
            "--delivery",
            str(delivery),
            "--builder",
            "python_pptx",
            "--output-dir",
            str(output),
            "--build-manifest",
            str(manifest),
        ],
        cwd=ROOT,
        text=True,
        capture_output=True,
        check=False,
    )
    assert result.returncode == 0, result.stderr
    assert (output / "deck.pptx").exists()
    data = json.loads(manifest.read_text(encoding="utf-8"))
    Draft202012Validator(json.loads((ROOT / "schemas" / "build-manifest.schema.json").read_text(encoding="utf-8"))).validate(data)
    assert data["status"] == "created"
    assert data["outputs"]["deck"] == "deck.pptx"


def test_python_adapter_preserves_and_applies_supported_style_contract(tmp_path: Path) -> None:
    from pptx import Presentation

    style = json.loads(STYLE.read_text(encoding="utf-8"))
    style["colors"].update({"background": "#123456", "primary": "#654321", "text_primary": "#ABCDEF", "surface_1": "#FEDCBA", "border": "#0A0B0C"})
    style["typography"]["font_primary"] = ["Courier New", "sans-serif"]
    style["typography"]["title_sizes_pt"]["slide"] = 31
    style["grid"]["margin_left_in"] = 0.91
    style["spacing"]["scale"]["md"] = 0.27
    style["footer_tokens"]["enabled"] = True
    plan = PythonPptxAdapter().plan(sample_ppt_ir(), style, sample_delivery_plan())
    assert plan.style_contract == style
    assert any("footer_tokens.enabled" in warning for warning in plan.warnings)

    result = PythonPptxAdapter().build(plan, tmp_path)
    assert any("footer_tokens.enabled" in warning for warning in result.warnings)
    deck = Presentation(result.pptx)
    assert round(deck.slide_width / 914400, 3) == 13.333
    assert round(deck.slide_height / 914400, 3) == 7.5
    slide = deck.slides[1]
    assert str(slide.background.fill.fore_color.rgb) == "123456"
    title = next(shape for shape in slide.shapes if shape.has_text_frame and shape.text == "Text and tables are native")
    run = title.text_frame.paragraphs[0].runs[0]
    assert run.font.name == "Courier New"
    assert round(run.font.size.pt) == 31
    assert str(run.font.color.rgb) == "654321"
    assert round(title.left / 914400, 2) == 0.91
    table = next(shape.table for shape in slide.shapes if shape.has_table)
    assert str(table.cell(0, 0).fill.fore_color.rgb) == "654321"
    assert str(table.cell(1, 0).fill.fore_color.rgb) == "123456"
    assert table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.name == "Courier New"


def test_python_adapter_warns_for_each_unconsumed_style_path_without_false_positives() -> None:
    style = json.loads(EDITORIAL_STYLE.read_text(encoding="utf-8"))

    plan = PythonPptxAdapter().plan(sample_ppt_ir(), style, sample_delivery_plan())
    warning_paths = {
        warning.rsplit(": ", 1)[-1]
        for warning in plan.warnings
        if warning.startswith("Unsupported style field ignored by PythonPptxAdapter: ")
    }

    assert {
        "colors.accent",
        "typography.font_editorial",
        "grid.margin_right_in",
        "spacing.rules.section_gap",
        "table_tokens.default.border_width_pt",
    } <= warning_paths
    assert {
        "colors.background",
        "colors.primary",
        "typography.font_primary",
        "typography.title_sizes_pt.slide",
        "typography.body_sizes_pt.normal",
        "typography.body_sizes_pt.small",
        "grid.margin_left_in",
        "grid.margin_top_in",
        "spacing.rules.card_gap",
        "spacing.scale.md",
    }.isdisjoint(warning_paths)


def test_python_adapter_preserves_native_required_chart_and_process_semantics(tmp_path: Path) -> None:
    from pptx import Presentation

    plan = sample_ppt_ir()
    plan["deck"]["logical_slide_count"] = 1
    plan["slides"] = [{
        **plan["slides"][1],
        "id": "S01",
        "index": 1,
        "title": "Native chart and process",
        "objects": [
            {
                "id": "chart-1", "type": "chart", "component_type": "bar_chart",
                "content": {"categories": ["Text", "Table"], "series": [{"name": "Count", "values": [4, 1]}]},
                "editability": "native_required", "source_refs": [],
            },
            {
                "id": "process-1", "type": "shape", "component_type": "process",
                "content": {"nodes": [{"label": "Plan"}, {"label": "Build"}, {"label": "Inspect"}]},
                "editability": "native_required", "source_refs": [],
            },
        ],
    }]
    delivery = sample_delivery_plan()
    delivery["slides"] = [{"slide_id": "S01", "objects": [
        {**delivery["slides"][0]["objects"][0], "slide_id": "S01", "object_id": "chart-1", "component_type": "bar_chart", "selected_route": "native_chart"},
        {**delivery["slides"][0]["objects"][0], "slide_id": "S01", "object_id": "process-1", "component_type": "process", "selected_route": "native_diagram"},
    ]}]
    result = PythonPptxAdapter().build(PythonPptxAdapter().plan(plan, json.loads(STYLE.read_text()), delivery), tmp_path)

    assert result.status == "created", result.errors
    assert result.fallbacks == []
    deck = Presentation(result.pptx)
    slide = deck.slides[0]
    chart_shape = next(shape for shape in slide.shapes if shape.has_chart)
    assert [str(category.label) for category in chart_shape.chart.plots[0].categories] == ["Text", "Table"]
    assert [(series.name, list(series.values)) for series in chart_shape.chart.series] == [("Count", [4.0, 1.0])]
    editable_text = {shape.text.strip() for shape in slide.shapes if shape.has_text_frame and shape.text.strip()}
    assert {"Plan", "Build", "Inspect"} <= editable_text
    assert {(item["object_id"], item["actual_route"]) for item in result.object_results} == {
        ("chart-1", "native_chart"), ("process-1", "native_diagram")
    }


def test_process_uses_single_bound_arrow_connectors_without_arrowhead_shapes(tmp_path: Path) -> None:
    from pptx import Presentation

    plan = sample_ppt_ir()
    plan["deck"]["logical_slide_count"] = 1
    plan["slides"] = [{
        **plan["slides"][1],
        "id": "S01",
        "index": 1,
        "title": "Bound process connectors",
        "objects": [{
            "id": "process-1", "type": "shape", "component_type": "process",
            "content": {"nodes": [{"label": "Plan"}, {"label": "Build"}, {"label": "Inspect"}]},
            "editability": "native_required", "source_refs": [],
        }],
    }]
    delivery = sample_delivery_plan()
    delivery["slides"] = [{"slide_id": "S01", "objects": [
        {**delivery["slides"][0]["objects"][0], "slide_id": "S01", "object_id": "process-1", "component_type": "process", "selected_route": "native_diagram"},
    ]}]

    result = PythonPptxAdapter().build(PythonPptxAdapter().plan(plan, json.loads(STYLE.read_text()), delivery), tmp_path)
    deck = Presentation(result.pptx)
    slide = deck.slides[0]
    connectors = [shape for shape in slide.shapes if shape.shape_type == 9]
    assert len(connectors) == 2
    assert not any(shape.auto_shape_type == 7 for shape in slide.shapes if shape.shape_type == 1)
    for connector in connectors:
        xml = connector._element.xml
        assert "<a:stCxn" in xml
        assert "<a:endCxn" in xml
        assert '<a:tailEnd type="triangle"' in xml


def test_cli_builder_override_is_truthful_and_preserves_plan_metadata(tmp_path: Path) -> None:
    ppt_ir = tmp_path / "ppt-ir.json"; style = tmp_path / "style.json"; delivery = tmp_path / "delivery.json"
    manifest = tmp_path / "manifest.json"; output = tmp_path / "out"
    plan = sample_delivery_plan()
    plan["builder"]["requested"] = "pptxgenjs"
    plan["builder"]["selected"] = "pptxgenjs"
    write_json(ppt_ir, sample_ppt_ir()); style.write_text(STYLE.read_text()); write_json(delivery, plan)
    completed = subprocess.run([sys.executable, str(BUILD), "--ppt-ir", str(ppt_ir), "--style", str(style), "--delivery", str(delivery), "--builder", "python_pptx", "--output-dir", str(output), "--build-manifest", str(manifest)], cwd=ROOT, text=True, capture_output=True)
    assert completed.returncode == 0, completed.stderr
    data = json.loads(manifest.read_text())
    assert data["builder"]["selected"] == "python_pptx"
    assert data["builder"]["requested"] == "pptxgenjs"
    assert "EXPLICIT_BUILDER_OVERRIDE" in data["builder"]["selection_reasons"]
    assert any("planned builder pptxgenjs" in warning for warning in data["warnings"])
