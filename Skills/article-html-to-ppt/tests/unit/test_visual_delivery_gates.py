from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[2]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from builders.selector import select_builder
from ppt_qa.package_inspector import inspect_package
from ppt_qa.slide_inspector import inspect_slides


def test_standard_rejects_minimal_python_builder_for_relationship_visual() -> None:
    ppt_ir = {
        "slides": [{
            "id": "S01",
            "primary_expression": "relationship_visual",
            "objects": [{"id": "architecture", "component_type": "layered_architecture", "editability": "native_required"}],
        }]
    }
    registry = {"components": [{
        "component_type": "layered_architecture",
        "preferred_delivery_route": "native_diagram",
        "builder_support": {"python_pptx": {"level": "full"}},
    }]}
    capabilities = {"builders": {"python_pptx": {"available": True, "version": "1", "components": {"layered_architecture": "full"}}}, "renderers": {"libreoffice": {"available": True}}}

    result = select_builder(ppt_ir, None, registry, capabilities, "standard", "python_pptx")
    assert "BUILDER_STANDARD_RELATIONSHIP_VISUAL_UNQUALIFIED" in result.errors


def test_inspector_rejects_unlabeled_solid_color_block(tmp_path: Path) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.util import Inches

    path = tmp_path / "orphan-block.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    block = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(2), Inches(2), Inches(2), Inches(1))
    block.fill.solid(); block.fill.fore_color.rgb = RGBColor(244, 96, 62)
    block.line.fill.background()
    deck.save(path)

    result = inspect_slides(path, inspect_package(path))
    issue_codes = {issue.issue_code for slide_result in result.slides for issue in slide_result.issues}
    assert "PPTX_ORPHAN_SOLID_COLOR_BLOCK" in issue_codes


def test_inspector_allows_explicit_material_layer(tmp_path: Path) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.util import Inches

    path = tmp_path / "material-block.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    block = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(2), Inches(2), Inches(2), Inches(1))
    block.name = "Material:card:summary"
    block.fill.solid(); block.fill.fore_color.rgb = RGBColor(244, 96, 62)
    deck.save(path)

    result = inspect_slides(path, inspect_package(path))
    issue_codes = {issue.issue_code for slide_result in result.slides for issue in slide_result.issues}
    assert "PPTX_ORPHAN_SOLID_COLOR_BLOCK" not in issue_codes
