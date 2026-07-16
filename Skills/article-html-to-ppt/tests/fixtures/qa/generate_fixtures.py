from __future__ import annotations

import copy
import json
import shutil
import zipfile
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parent


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


STYLE: dict[str, Any] = {
    "schema_version": "2.0",
    "style_id": "qa-style",
    "display_name": "QA Style",
    "audience": ["internal"],
    "formality": "professional",
    "presentation_context": ["review"],
    "aspect_ratios": ["16:9"],
    "colors": {
        "primary": "#1F4E79",
        "accent": "#E85D3F",
        "background": "#FFFFFF",
        "surface_1": "#F3F6FA",
        "surface_2": "#E7EEF7",
        "text_primary": "#1C1C1C",
        "text_secondary": "#555555",
        "border": "#C7D0DA",
        "positive": "#2E7D32",
        "warning": "#F9A825",
        "negative": "#C62828",
        "data_series": ["#1F4E79", "#E85D3F", "#2E7D32"],
        "allowed_opacity": [1.0, 0.8, 0.5]
    },
    "typography": {
        "font_primary": ["Arial", "Calibri", "sans-serif"],
        "font_editorial": ["Arial", "serif"],
        "font_mono": ["Courier New", "monospace"],
        "title_sizes_pt": {"h1": 30, "h2": 24},
        "body_sizes_pt": {"body": 18, "small": 12},
        "metric_sizes_pt": {"large": 36},
        "minimum_body_size_pt": 10,
        "minimum_footnote_size_pt": 8,
        "line_height": {"body": 1.2},
        "paragraph_spacing_pt": {"body": 6},
        "letter_spacing_pt": {"normal": 0},
        "weights": {"regular": 400, "bold": 700}
    },
    "grid": {
        "columns": 12,
        "rows": 8,
        "margin_left_in": 0.5,
        "margin_right_in": 0.5,
        "margin_top_in": 0.4,
        "margin_bottom_in": 0.4,
        "gutter_horizontal_in": 0.1,
        "gutter_vertical_in": 0.1,
        "title_zone_height_in": 0.8,
        "footer_zone_height_in": 0.2,
        "safe_zone_in": {"left": 0.1, "right": 0.1, "top": 0.1, "bottom": 0.1}
    },
    "spacing": {"unit": "in", "scale": {"sm": 0.1, "md": 0.2}, "rules": {"stack": "md"}},
    "shape_tokens": {},
    "shadow_tokens": {},
    "card_tokens": {
        "default": {"fill": "surface_1", "text_color": "text_primary", "border_color": "border"},
        "highlight": {"fill": "surface_2", "text_color": "text_primary", "border_color": "accent"},
        "metric": {"fill": "surface_1", "text_color": "text_primary", "number_color": "primary"},
        "risk": {"fill": "surface_1", "text_color": "text_primary", "border_color": "warning"},
        "quote": {"fill": "surface_1", "text_color": "text_secondary", "border_color": "border"},
        "comparison": {"fill": "surface_1", "text_color": "text_primary", "border_color": "border"},
        "source": {"fill": "surface_1", "text_color": "text_secondary", "border_color": "border"}
    },
    "table_tokens": {},
    "chart_tokens": {},
    "diagram_tokens": {},
    "image_tokens": {},
    "icon_tokens": {},
    "footer_tokens": {},
    "density_limits": {},
    "allowed_effects": [],
    "forbidden_drift": ["unregistered_hex"],
    "compatibility_aliases": []
}


def base_ppt_ir(slide_count: int = 1, primary_expression: str = "textual_argument") -> dict[str, Any]:
    slides = []
    for idx in range(1, slide_count + 1):
        slides.append({
            "id": f"S{idx:02d}",
            "index": idx,
            "slide_role": "data" if idx > 1 else "judgment",
            "title_role": "judgment",
            "title": f"QA Slide {idx}",
            "judgment": f"QA Slide {idx} stays editable",
            "message": f"Native body text for slide {idx}",
            "audience_question": "Can this deck stay editable?",
            "source_refs": [{"source_id": "src-1", "locator": "fixture", "claim_type": "direct"}],
            "primary_expression": primary_expression,
            "supporting_expressions": ["source_note"],
            "primary_anchor": "body",
            "objects": [{
                "id": f"text-{idx}",
                "type": "text",
                "semantic_role": "body",
                "content": f"Native body text for slide {idx}",
                "source_refs": [{"source_id": "src-1", "locator": "fixture", "claim_type": "direct"}],
                "editability": "native_required"
            }],
            "delivery_contract": {
                "preferred_route": "native_ppt",
                "editable_core": ["title", "body"],
                "raster_allowance": [],
                "forbidden_raster": ["core_text"]
            },
            "qa_expectations": ["native_text"]
        })
    return {
        "schema_version": "2.0",
        "deck": {
            "id": "qa-fixture",
            "title": "QA Fixture",
            "source_type": "article",
            "audience": "internal",
            "purpose": "qa",
            "language": "en",
            "aspect_ratio": "16:9",
            "production_profile": "standard",
            "logical_slide_count": slide_count
        },
        "sources": [{"source_id": "src-1", "type": "fixture", "title": "Fixture Source"}],
        "style_contract_ref": "style-contract.json",
        "asset_manifest_ref": "asset-manifest.json",
        "slides": slides
    }


def base_delivery(slide_count: int = 1, selected_route: str = "native_ppt", component_type: str = "body_argument") -> dict[str, Any]:
    slides = []
    for idx in range(1, slide_count + 1):
        slides.append({
            "slide_id": f"S{idx:02d}",
            "objects": [{
                "slide_id": f"S{idx:02d}",
                "object_id": f"text-{idx}",
                "component_type": component_type,
                "preferred_route": selected_route,
                "selected_route": selected_route,
                "decision": "selected",
                "reason_codes": [],
                "editable_core": ["title", "body"],
                "rasterized_parts": [],
                "svg_parts": [],
                "native_overlay_parts": [],
                "qa_checks": ["native_text"],
                "fallback_chain": []
            }]
        })
    return {
        "schema_version": "1.0",
        "ppt_ir_ref": "ppt-ir.json",
        "style_contract_ref": "style-contract.json",
        "component_registry_ref": "component-registry.json",
        "profile": "standard",
        "builder": {"requested": "fixture", "selected": "fixture"},
        "slides": slides,
        "summary": {"total_objects": slide_count, "fallback_count": 0, "unsupported_count": 0, "risk_codes": []}
    }


def set_text(shape: Any, text: str, size: int = 18, color: str = "#1C1C1C") -> None:
    shape.text = text
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(size)
            run.font.color.rgb = rgb(color)


def blank_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    return prs


def add_title_body(slide: Any, title: str = "QA Slide 1", body: str = "Native body text for slide 1") -> None:
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(5.8), Inches(0.55))
    set_text(title_box, title, 28)
    body_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(5.5), Inches(0.8))
    set_text(body_box, body, 18)


def write_contracts(path: Path, ppt_ir: dict[str, Any], delivery: dict[str, Any], expected: list[str]) -> None:
    (path / "ppt-ir.json").write_text(json.dumps(ppt_ir, indent=2) + "\n", encoding="utf-8")
    (path / "style-contract.json").write_text(json.dumps(STYLE, indent=2) + "\n", encoding="utf-8")
    (path / "delivery-plan.json").write_text(json.dumps(delivery, indent=2) + "\n", encoding="utf-8")
    (path / "expected-issues.json").write_text(json.dumps(expected, indent=2) + "\n", encoding="utf-8")


def make_image(path: Path, color: str = "#D9E8F7", text: str = "raster") -> None:
    image = Image.new("RGB", (1600, 900), color)
    draw = ImageDraw.Draw(image)
    draw.rectangle((80, 80, 1520, 820), outline="#1F4E79", width=8)
    draw.text((150, 410), text, fill="#1C1C1C")
    image.save(path)


def save_valid(path: Path) -> None:
    prs = blank_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_body(slide)
    table = slide.shapes.add_table(3, 3, Inches(0.7), Inches(2.2), Inches(4.0), Inches(1.2)).table
    for row in range(3):
        for col in range(3):
            table.cell(row, col).text = f"{row + 1}-{col + 1}"
    chart_data = CategoryChartData()
    chart_data.categories = ["A", "B", "C"]
    chart_data.add_series("Series", (4, 7, 5))
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(5.4), Inches(2.0), Inches(4.0), Inches(2.4), chart_data)
    prs.save(path / "deck.pptx")
    ppt_ir = base_ppt_ir()
    ppt_ir["slides"][0]["primary_expression"] = "data_visual"
    delivery = base_delivery(selected_route="native_chart", component_type="bar_chart")
    write_contracts(path, ppt_ir, delivery, [])


def save_whole_slide_image(path: Path, name: str = "whole-slide-image", delivery_route: str = "raster_component", overlay: bool = False) -> None:
    img = path / "image.png"
    make_image(img, text=name)
    prs = blank_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pic = slide.shapes.add_picture(str(img), Inches(0.0), Inches(0.0), Inches(13.333333), Inches(7.5))
    pic.name = "Background Image" if overlay else "Full Slide Raster"
    if overlay:
        add_title_body(slide)
    prs.save(path / "deck.pptx")
    ppt_ir = base_ppt_ir()
    ppt_ir["slides"][0]["objects"][0].update({
        "id": "image-1",
        "type": "image",
        "component_type": "qa_raster_image",
        "semantic_role": "visual",
        "content": "Raster fixture image",
        "editability": "raster_allowed",
    })
    ppt_ir["slides"][0]["delivery_contract"]["preferred_route"] = delivery_route
    ppt_ir["slides"][0]["delivery_contract"]["raster_allowance"] = ["image-1"]
    delivery = base_delivery(selected_route=delivery_route, component_type="qa_raster_image")
    delivery["slides"][0]["objects"][0]["object_id"] = "image-1"
    delivery["slides"][0]["objects"][0]["editable_core"] = []
    if delivery_route == "background_image":
        delivery["slides"][0]["objects"][0]["rasterized_parts"] = ["background"]
    write_contracts(path, ppt_ir, delivery, [] if overlay else ["PPTX_WHOLE_SLIDE_RASTER"])
    img.unlink()


def rewrite_without_member(pptx: Path, predicate: Any) -> None:
    tmp = pptx.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(pptx, "r") as source, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as target:
        for item in source.infolist():
            if predicate(item.filename):
                continue
            target.writestr(item, source.read(item.filename))
    tmp.replace(pptx)


def rewrite_slide_rels_target(pptx: Path) -> None:
    tmp = pptx.with_suffix(".tmp.pptx")
    rels = "ppt/slides/_rels/slide1.xml.rels"
    with zipfile.ZipFile(pptx, "r") as source, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as target:
        for item in source.infolist():
            data = source.read(item.filename)
            if item.filename == rels:
                text = data.decode("utf-8").replace("../media/image1.png", "../media/missing-image.png")
                data = text.encode("utf-8")
            target.writestr(item, data)
    tmp.replace(pptx)


def save_picture_component(path: Path, route: str, issue: str) -> None:
    save_whole_slide_image(path, name=issue, delivery_route=route)
    delivery = json.loads((path / "delivery-plan.json").read_text(encoding="utf-8"))
    delivery["slides"][0]["objects"][0]["selected_route"] = route
    delivery["slides"][0]["objects"][0]["preferred_route"] = route
    (path / "delivery-plan.json").write_text(json.dumps(delivery, indent=2) + "\n", encoding="utf-8")
    (path / "expected-issues.json").write_text(json.dumps([issue], indent=2) + "\n", encoding="utf-8")


def save_text_overflow(path: Path) -> None:
    prs = blank_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_body(slide)
    box = slide.shapes.add_textbox(Inches(6), Inches(1), Inches(1.0), Inches(0.25))
    set_text(box, "This is a very long text block that should overflow the tiny box by estimation.", 18)
    prs.save(path / "deck.pptx")
    write_contracts(path, base_ppt_ir(), base_delivery(), ["TEXT_OVERFLOW_RISK"])


def save_out_of_bounds(path: Path) -> None:
    prs = blank_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_body(slide)
    box = slide.shapes.add_textbox(Inches(12.8), Inches(7.1), Inches(1.2), Inches(0.5))
    set_text(box, "Off", 18)
    prs.save(path / "deck.pptx")
    write_contracts(path, base_ppt_ir(), base_delivery(), ["PPTX_TEXT_OUT_OF_BOUNDS"])


def save_color_drift(path: Path) -> None:
    prs = blank_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_body(slide)
    shape = slide.shapes.add_shape(1, Inches(6), Inches(1.5), Inches(2), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("#FF00FF")
    shape.line.color.rgb = rgb("#FF00FF")
    prs.save(path / "deck.pptx")
    write_contracts(path, base_ppt_ir(), base_delivery(), ["STYLE_COLOR_DRIFT"])


def save_font_too_small(path: Path) -> None:
    prs = blank_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_body(slide)
    box = slide.shapes.add_textbox(Inches(6), Inches(1), Inches(2), Inches(0.4))
    set_text(box, "Too small", 6)
    prs.save(path / "deck.pptx")
    write_contracts(path, base_ppt_ir(), base_delivery(), ["PPTX_FONT_SIZE_BELOW_MINIMUM"])


def save_blank(path: Path) -> None:
    prs = blank_prs()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(path / "deck.pptx")
    write_contracts(path, base_ppt_ir(), base_delivery(), ["PPTX_BLANK_SLIDE"])


def save_fragmented(path: Path) -> None:
    prs = blank_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for idx in range(155):
        x = 0.1 + (idx % 31) * 0.18
        y = 0.1 + (idx // 31) * 0.18
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(0.03), Inches(0.03))
        set_text(box, ".", 5)
    prs.save(path / "deck.pptx")
    write_contracts(path, base_ppt_ir(), base_delivery(), ["PPTX_TINY_OBJECT_OVERLOAD"])


def save_complex_diagram(path: Path) -> None:
    prs = blank_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_body(slide, body="Native body text for slide 1")
    for idx in range(6):
        shape = slide.shapes.add_shape(1, Inches(0.8 + idx * 1.4), Inches(2.3), Inches(1), Inches(0.5))
        shape.text = f"N{idx + 1}"
        if idx:
            slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.8 + (idx - 1) * 1.4 + 1), Inches(2.55), Inches(0.8 + idx * 1.4), Inches(2.55))
    prs.save(path / "deck.pptx")
    write_contracts(path, base_ppt_ir(), base_delivery(selected_route="native_diagram", component_type="process_diagram"), [])


def main() -> None:
    if ROOT.exists():
        for item in ROOT.iterdir():
            if item.name == "generate_fixtures.py":
                continue
            if item.is_dir():
                shutil.rmtree(item)
    fixtures = {
        "valid-native-deck": save_valid,
        "whole-slide-image": save_whole_slide_image,
        "text-overflow": save_text_overflow,
        "object-out-of-bounds": save_out_of_bounds,
        "color-drift": save_color_drift,
        "font-too-small": save_font_too_small,
        "blank-slide": save_blank,
        "fragmented-svg-conversion": save_fragmented,
        "complex-diagram": save_complex_diagram,
    }
    for name, func in fixtures.items():
        path = ROOT / name
        path.mkdir(parents=True, exist_ok=True)
        func(path)
    for name, route, issue in [
        ("rasterized-table", "native_table", "PPTX_TABLE_NOT_NATIVE"),
        ("rasterized-chart", "native_chart", "PPTX_CHART_NOT_NATIVE"),
        ("route-deviation", "native_chart", "PPTX_RASTER_ROUTE_UNDECLARED"),
    ]:
        path = ROOT / name
        path.mkdir(parents=True, exist_ok=True)
        save_picture_component(path, route, issue)
    for name, rewrite, expected in [
        ("missing-media", lambda deck: rewrite_without_member(deck, lambda filename: filename.startswith("ppt/media/")), ["PPTX_MISSING_MEDIA"]),
        ("broken-relationship", rewrite_slide_rels_target, ["PPTX_MISSING_MEDIA"]),
    ]:
        path = ROOT / name
        path.mkdir(parents=True, exist_ok=True)
        save_whole_slide_image(path, name=name)
        rewrite(path / "deck.pptx")
        (path / "expected-issues.json").write_text(json.dumps(expected, indent=2) + "\n", encoding="utf-8")


if __name__ == "__main__":
    main()
