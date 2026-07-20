from __future__ import annotations

import json
import subprocess
import sys
import zipfile
from collections import Counter

from pptx import Presentation
from pathlib import Path

import pytest
from jsonschema import Draft202012Validator

ROOT = Path(__file__).resolve().parents[2]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from ppt_qa.package_inspector import inspect_package

FIXTURES = ROOT / "tests" / "fixtures" / "v2-acceptance"
BUILD = ROOT / "scripts" / "build_deck.py"
SCHEMAS = ROOT / "schemas"

SHARED_STYLE = "style-contract-editorial.json"
SHARED_DELIVERY = "delivery-plan-shared.json"
BUILDERS = ["python_pptx", "pptxgenjs"]


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def validate(document: dict, schema_name: str) -> None:
    schema = load_json(SCHEMAS / schema_name)
    Draft202012Validator(schema).validate(document)


@pytest.fixture(scope="module")
def acceptance_ir() -> dict:
    document = load_json(FIXTURES / "ppt-ir.json")
    validate(document, "ppt-ir.schema.json")
    assert len(document["slides"]) == document["deck"]["logical_slide_count"] == 9
    assert [slide["slide_role"] for slide in document["slides"]] == [
        "cover", "judgment", "data", "data", "data", "diagram", "judgment", "instruction", "closing"
    ]
    return document


def test_acceptance_contracts_are_schema_valid(acceptance_ir: dict) -> None:
    style = load_json(FIXTURES / SHARED_STYLE)
    delivery = load_json(FIXTURES / SHARED_DELIVERY)
    validate(style, "style-contract.schema.json")
    validate(delivery, "delivery-plan.schema.json")
    # Brief-required compatibility fixtures are aliases of the canonical plan, preventing drift.
    assert (FIXTURES / "delivery-plan-python.json").resolve() == (FIXTURES / SHARED_DELIVERY).resolve()
    assert (FIXTURES / "delivery-plan-pptxgenjs.json").resolve() == (FIXTURES / SHARED_DELIVERY).resolve()
    validate(load_json(FIXTURES / "style-contract-technical.json"), "style-contract.schema.json")
    assert delivery["style_contract_ref"] == SHARED_STYLE
    assert delivery["builder"]["selected"] == "unknown"
    assert delivery["summary"]["total_objects"] == sum(len(slide.get("objects", [])) for slide in acceptance_ir["slides"])


def build_case(tmp_path: Path, builder: str) -> dict:
    output_dir = tmp_path / builder
    manifest_path = output_dir / "build-manifest.json"
    completed = subprocess.run(
        [
            sys.executable,
            str(BUILD),
            "--ppt-ir", str(FIXTURES / "ppt-ir.json"),
            "--style", str(FIXTURES / SHARED_STYLE),
            "--delivery", str(FIXTURES / SHARED_DELIVERY),
            "--builder", builder,
            "--output-dir", str(output_dir),
            "--build-manifest", str(manifest_path),
            "--json-output",
        ],
        cwd=ROOT,
        text=True,
        capture_output=True,
        check=False,
    )
    assert completed.returncode == 0, completed.stderr or completed.stdout
    payload = json.loads(completed.stdout)
    payload["manifest_path"] = manifest_path
    payload["deck_path"] = output_dir / payload["build_manifest"]["outputs"]["deck"]
    return payload


def _object_strings(obj: dict) -> set[str]:
    content = obj.get("content")
    values: set[str] = set()
    def visit(value):
        if isinstance(value, str) and value.strip(): values.add(value.strip())
        elif isinstance(value, dict):
            for nested in value.values(): visit(nested)
        elif isinstance(value, list):
            for nested in value: visit(nested)
    visit(content)
    return values


def _chart_semantics(shape) -> tuple[list[str], list[tuple[str, list[float]]]]:
    chart = shape.chart
    categories = [str(category.label) for category in chart.plots[0].categories]
    series = [(item.name, list(item.values)) for item in chart.series]
    return categories, series


def _assert_native_package_evidence(deck_path: Path, acceptance_ir: dict, style: dict, builder: str) -> None:
    deck = Presentation(deck_path)
    assert len(deck.slides) == len(acceptance_ir["slides"])
    for slide_ir, slide in zip(acceptance_ir["slides"], deck.slides):
        native_text = {shape.text.strip() for shape in slide.shapes if shape.has_text_frame and shape.text.strip()}
        assert slide_ir["title"] in native_text, f"missing editable native title on {slide_ir['id']}"
        title_shape = next(shape for shape in slide.shapes if shape.has_text_frame and shape.text.strip() == slide_ir["title"])
        title_run = title_shape.text_frame.paragraphs[0].runs[0]
        expected_title_font = ((style["typography"].get("font_editorial") or style["typography"]["font_primary"])[0] if builder == "pptxgenjs" else style["typography"]["font_primary"][0])
        assert title_run.font.name == expected_title_font
        title_color = style["colors"]["text_primary" if builder == "pptxgenjs" else "primary"]
        assert str(title_run.font.color.rgb) == title_color.lstrip("#").upper()
        assert str(slide.background.fill.fore_color.rgb) == style["colors"]["background"].lstrip("#").upper()
        assert any(shape.has_text_frame for shape in slide.shapes), f"missing editable native shape on {slide_ir['id']}"
        all_text = "\n".join(shape.text for shape in slide.shapes if shape.has_text_frame)
        native_text_parts = {part.strip() for part in all_text.split("\n") if part.strip()}
        for obj in slide_ir.get("objects", []):
            if obj.get("editability") != "native_required":
                continue
            if obj.get("type") == "table":
                tables = [shape.table for shape in slide.shapes if shape.has_table]
                assert len(tables) == 1, f"expected one native table on {slide_ir['id']}"
                table_text = {cell.text for row in tables[0].rows for cell in row.cells}
                assert _object_strings(obj) <= table_text
            elif obj.get("type") == "chart":
                charts = [shape for shape in slide.shapes if shape.has_chart]
                assert len(charts) == 1, f"expected one native editable chart on {slide_ir['id']}"
                content = obj["content"]
                expected_series = [(series["name"], series["values"]) for series in content["series"]]
                assert _chart_semantics(charts[0]) == (content["categories"], expected_series)
            elif obj.get("component_type") == "process":
                labels = {node["label"] for node in obj["content"]["nodes"]}
                assert labels <= native_text_parts, f"missing editable process labels for {obj['id']}"
            else:
                assert _object_strings(obj) & native_text_parts, f"missing IR native text for {obj['id']}"


def test_shared_ir_builds_comparable_real_packages(tmp_path: Path, acceptance_ir: dict) -> None:
    builds = {builder: build_case(tmp_path, builder) for builder in BUILDERS}
    shared_style = load_json(FIXTURES / SHARED_STYLE)
    shared_delivery = load_json(FIXTURES / SHARED_DELIVERY)
    declared_pairs = {
        (slide["id"], obj["id"])
        for slide in acceptance_ir["slides"]
        for obj in slide.get("objects", [])
    }
    delivery_pairs = {
        (slide["slide_id"], obj["object_id"])
        for slide in shared_delivery["slides"]
        for obj in slide.get("objects", [])
    }
    assert delivery_pairs == declared_pairs
    declared_count = len(declared_pairs)
    assert shared_delivery["summary"]["total_objects"] == declared_count
    slide_counts = set()

    for builder, payload in builds.items():
        manifest = payload["build_manifest"]
        result = payload["result"]
        deck_path = payload["deck_path"]
        validate(manifest, "build-manifest.schema.json")

        assert deck_path.is_file() and deck_path.stat().st_size > 0
        assert zipfile.is_zipfile(deck_path)
        _assert_native_package_evidence(deck_path, acceptance_ir, shared_style, builder)
        inspection = inspect_package(deck_path, ppt_ir=acceptance_ir, build_manifest=manifest)
        assert inspection.status == "passed", [issue.__dict__ for issue in inspection.issues]
        slide_counts.add(inspection.slide_count)

        assert manifest["status"] == "created"
        assert manifest["builder"]["selected"] == builder
        assert manifest["builder_profile"] == builder
        assert result["builder"] == builder
        assert load_json(FIXTURES / SHARED_STYLE) == shared_style
        assert load_json(FIXTURES / SHARED_DELIVERY) == shared_delivery
        assert manifest["metrics"]["slide_count"] == inspection.slide_count
        assert manifest["metrics"]["object_count"] == len(result["object_results"]) == declared_count
        assert manifest["metrics"]["fallback_count"] == len(result["fallbacks"])
        assert manifest.get("fallbacks", []) == result["fallbacks"]
        deviations = [item for item in result["object_results"] if item["planned_route"] != item["actual_route"]]
        expected_keys = Counter((item["slide_id"], item["object_id"], item["planned_route"], item["actual_route"]) for item in deviations)
        fallback_keys = Counter((item["slide_id"], item["object_id"], item["planned_route"], item["actual_route"]) for item in result["fallbacks"])
        assert fallback_keys == expected_keys
        ir_by_pair = {(slide["id"], obj["id"]): obj for slide in acceptance_ir["slides"] for obj in slide.get("objects", [])}
        for fallback in result["fallbacks"]:
            obj = ir_by_pair[(fallback["slide_id"], fallback["object_id"])]
            if obj["editability"] == "native_required":
                assert fallback["editable_core_preserved"] is True

        created_pairs = {
            (item["slide_id"], item["object_id"])
            for item in result["object_results"]
            if item["status"] == "created"
        }
        assert created_pairs == declared_pairs
        assert not any(
            error.get("code") == "UNSUPPORTED_NATIVE_REQUIRED_OBJECT"
            for error in result["errors"]
            if isinstance(error, dict)
        )

    assert slide_counts == {9}
